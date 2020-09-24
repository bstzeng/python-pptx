from pptx import Presentation
from pptx.util import Inches
from pptx.util import Cm
import copy
import six

def duplicate_slide(pres,index):
        template = pres.slides[index]
        try:
            blank_slide_layout = pres.slide_layouts[12]
        except:
            blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts)]

        copied_slide = pres.slides.add_slide(blank_slide_layout)

        for shp in template.shapes:
            el = shp.element
            newel = copy.deepcopy(el)
            copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

        for _, value in six.iteritems(template.part.rels):
            # Make sure we don't copy a notesSlide relation as that won't exist
            if "notesSlide" not in value.reltype:
                copied_slide.part.rels.add_relationship(value.reltype,
                                                value._target,
                                                value.rId)

        return copied_slide


img_path = '1.jpg'

prs = Presentation('PE1template.pptx')
blank_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = Cm(5)
#pic = slide.shapes.add_picture(img_path, left, top)

height = Cm(5.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

img_path = '2.jpg'

slide = prs.slides.add_slide(blank_slide_layout)

left = top = Inches(1)
#pic = slide.shapes.add_picture(img_path, left, top)

left = Inches(5)
height = Inches(5.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

pres=Presentation('batch summary.pptx')
copied_slide = duplicate_slide(pres, 1)

prs.save('test.pptx')